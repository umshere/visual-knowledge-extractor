from __future__ import annotations

import hashlib
from pathlib import Path
from typing import Any, Dict, List, Optional

from pptx import Presentation

from ocr_images import ocr_image_bytes_safe


def _sha1(b: bytes) -> str:
    return hashlib.sha1(b).hexdigest()


def extract_pptx(pptx_path: Path, out_dir: Path) -> Dict[str, Any]:
    """Extract slide text, notes, images, and OCR into a structured dict."""

    prs = Presentation(str(pptx_path))
    images_dir = out_dir / "images"
    images_dir.mkdir(parents=True, exist_ok=True)

    slides_out: List[Dict[str, Any]] = []

    for idx, slide in enumerate(prs.slides, start=1):
        slide_text_runs: List[str] = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False) and shape.text_frame is not None:
                txt = (shape.text or "").strip()
                if txt:
                    slide_text_runs.append(txt)

        notes_text = ""
        try:
            if slide.has_notes_slide and slide.notes_slide is not None:
                # notes text is spread across placeholders
                notes = []
                for shp in slide.notes_slide.shapes:
                    if getattr(shp, "has_text_frame", False) and shp.text_frame is not None:
                        t = (shp.text or "").strip()
                        if t:
                            notes.append(t)
                notes_text = "\n".join(notes).strip()
        except Exception:
            notes_text = ""

        extracted_images: List[Dict[str, Any]] = []
        for shape in slide.shapes:
            if shape.shape_type == 13:  # MSO_SHAPE_TYPE.PICTURE
                try:
                    blob = shape.image.blob
                    ext = shape.image.ext or "img"
                    digest = _sha1(blob)[:12]
                    img_name = f"slide-{idx:03d}-{digest}.{ext}"
                    img_path = images_dir / img_name
                    img_path.write_bytes(blob)

                    ocr_text = ocr_image_bytes_safe(blob)

                    extracted_images.append(
                        {
                            "file": str(img_path),
                            "ext": ext,
                            "sha1": _sha1(blob),
                            "ocr_text": ocr_text,
                        }
                    )
                except Exception as e:
                    extracted_images.append({"error": f"image_extract_failed: {e}"})

        slides_out.append(
            {
                "index": idx,
                "text": "\n".join(slide_text_runs).strip(),
                "notes": notes_text,
                "images": extracted_images,
            }
        )

    # Very light heuristics for a "knowledge doc" (no LLM here)
    all_text = "\n\n".join(
        [s.get("text", "") + "\n" + s.get("notes", "") + "\n" + "\n".join([i.get("ocr_text", "") for i in s.get("images", []) if isinstance(i, dict)])
         for s in slides_out]
    ).strip()

    return {
        "source_file": str(pptx_path),
        "slide_count": len(slides_out),
        "slides": slides_out,
        "summary": {
            "executive_summary": _exec_summary(all_text),
            "key_points": _key_points(all_text),
        },
    }


def _exec_summary(text: str) -> str:
    if not text:
        return ""
    # naive summary: first ~800 chars of normalized text
    t = " ".join(text.split())
    return t[:800] + ("…" if len(t) > 800 else "")


def _key_points(text: str) -> List[str]:
    if not text:
        return []
    # naive key points: top lines that look like bullets
    lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
    bullets = [ln for ln in lines if ln.startswith(("-", "•", "*"))]
    out = []
    for b in bullets:
        cleaned = b.lstrip("-*• ").strip()
        if cleaned and cleaned not in out:
            out.append(cleaned)
        if len(out) >= 12:
            break
    return out
