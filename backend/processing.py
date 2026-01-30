from __future__ import annotations

import logging
import os
from pathlib import Path
from typing import Callable, List, Optional

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from PIL import Image

try:
    import pytesseract
except Exception:  # pragma: no cover - optional dependency
    pytesseract = None

from .models import ExtractedImageInfo, KnowledgeDoc, SlideInfo

logger = logging.getLogger(__name__)


def _tesseract_available() -> bool:
    if pytesseract is None:
        return False
    try:
        _ = pytesseract.get_tesseract_version()
        return True
    except Exception:
        return False


def _extract_text_items(slide) -> List[str]:
    items: List[str] = []
    for shape in slide.shapes:
        if not hasattr(shape, "text"):
            continue
        text = (shape.text or "").strip()
        if text:
            items.append(text)
    return items


def _extract_notes(slide) -> Optional[str]:
    if not slide.has_notes_slide:
        return None
    notes_text_frame = slide.notes_slide.notes_text_frame
    text = (notes_text_frame.text or "").strip()
    return text or None


def _ocr_image(image_path: Path) -> Optional[str]:
    if not _tesseract_available():
        return None
    try:
        with Image.open(image_path) as image:
            return pytesseract.image_to_string(image).strip() or None
    except Exception as exc:
        logger.warning("OCR failed for %s: %s", image_path, exc)
        return None


def extract_knowledge_doc(
    pptx_path: Path,
    job_dir: Path,
    on_slide: Optional[Callable[[List[SlideInfo], int, int], None]] = None,
) -> KnowledgeDoc:
    presentation = Presentation(str(pptx_path))
    slides: List[SlideInfo] = []

    images_dir = job_dir / "images"
    images_dir.mkdir(parents=True, exist_ok=True)

    ocr_enabled = _tesseract_available()
    if not ocr_enabled:
        logger.info("Tesseract not available; OCR will be skipped.")

    total_slides = len(presentation.slides)
    for slide_index, slide in enumerate(presentation.slides, start=1):
        title = None
        if slide.shapes.title is not None:
            title = (slide.shapes.title.text or "").strip() or None

        text_items = _extract_text_items(slide)
        notes = _extract_notes(slide)

        images: List[ExtractedImageInfo] = []
        image_counter = 0

        for shape in slide.shapes:
            if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                continue

            image_counter += 1
            image = shape.image
            ext = image.ext
            image_id = f"slide{slide_index}_image{image_counter}"
            filename = f"{image_id}.{ext}"
            image_path = images_dir / filename
            with image_path.open("wb") as f:
                f.write(image.blob)

            ocr_text = _ocr_image(image_path) if ocr_enabled else None
            images.append(
                ExtractedImageInfo(
                    image_id=image_id,
                    filename=str(image_path),
                    slide_index=slide_index,
                    ocr_text=ocr_text,
                )
            )

        slides.append(
            SlideInfo(
                slide_index=slide_index,
                title=title,
                text_items=text_items,
                notes=notes,
                images=images,
            )
        )
        if on_slide:
            on_slide(slides, slide_index, total_slides)

    return KnowledgeDoc(
        source_filename=pptx_path.name,
        slide_count=len(slides),
        slides=slides,
    )


def generate_markdown(doc: KnowledgeDoc) -> str:
    lines: List[str] = []
    lines.append(f"# Knowledge Document: {doc.source_filename}")
    lines.append("")
    for slide in doc.slides:
        lines.append(f"## Slide {slide.slide_index}")
        if slide.title:
            lines.append(f"**Title:** {slide.title}")
        if slide.text_items:
            lines.append("")
            lines.append("**Text:**")
            for item in slide.text_items:
                lines.append(f"- {item}")
        if slide.notes:
            lines.append("")
            lines.append("**Speaker Notes:**")
            lines.append(slide.notes)
        if slide.images:
            lines.append("")
            lines.append("**Images:**")
            for image in slide.images:
                lines.append(f"- {os.path.basename(image.filename)}")
                if image.ocr_text:
                    lines.append(f"  - OCR: {image.ocr_text}")
        lines.append("")
    return "\n".join(lines).strip() + "\n"
